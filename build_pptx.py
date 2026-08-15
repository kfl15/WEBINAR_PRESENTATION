from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUT = "developer_side_memory_optimization.pptx"

W = Inches(13.333)
H = Inches(7.5)

INK = RGBColor(23, 32, 42)
MUTED = RGBColor(95, 109, 118)
PAPER = RGBColor(247, 245, 238)
SOFT_TEAL = RGBColor(229, 241, 240)
PANEL = RGBColor(255, 255, 255)
LINE = RGBColor(217, 212, 198)
TEAL = RGBColor(0, 109, 119)
GOLD = RGBColor(194, 124, 36)
RED = RGBColor(182, 70, 58)
GREEN = RGBColor(47, 125, 79)
BLUE = RGBColor(53, 92, 154)
PLUM = RGBColor(104, 65, 111)
DARK = RGBColor(23, 32, 42)
CODE = RGBColor(247, 245, 238)


def fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.color.rgb = LINE


def no_line(shape):
    shape.line.fill.background()


def add_text(slide, x, y, w, h, text, size=22, bold=False, color=INK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_multiline(slide, x, y, w, h, lines, size=21, color=INK, bullet=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for i, line in enumerate(lines):
      p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
      p.text = line
      p.level = 0
      if bullet:
          p.text = line
      for run in p.runs:
          run.font.name = "Aptos"
          run.font.size = Pt(size)
          run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, lines, size=19):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"- {line}"
        p.level = 0
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(size)
            run.font.color.rgb = INK
    return box


def add_card(slide, x, y, w, h, title, body=None, accent=TEAL, bullets=None, title_size=22, body_size=19):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shape, PANEL)
    shape.line.color.rgb = LINE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    fill(bar, accent)
    no_line(bar)
    add_text(slide, x + 0.18, y + 0.18, w - 0.32, 0.38, title, size=title_size, bold=True)
    if bullets:
        add_bullets(slide, x + 0.22, y + 0.68, w - 0.36, h - 0.82, bullets, size=body_size)
    elif body:
        add_text(slide, x + 0.18, y + 0.72, w - 0.32, h - 0.88, body, size=body_size, color=INK)
    return shape


def add_code(slide, x, y, w, h, text, size=15):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(shape, DARK)
    shape.line.color.rgb = DARK
    box = add_text(slide, x + 0.18, y + 0.14, w - 0.34, h - 0.24, text, size=size, color=CODE)
    for p in box.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = "Consolas"
    return shape


def add_tag(slide, x, y, text, w=None):
    if w is None:
        w = max(1.0, 0.13 * len(text) + 0.35)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.38))
    fill(shape, PANEL)
    add_text(slide, x + 0.05, y + 0.07, w - 0.1, 0.23, text, size=13, bold=True, align=PP_ALIGN.CENTER)
    return w


def base_slide(prs, kicker, title, number, total, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    fill(bg, PAPER)
    no_line(bg)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.9), H)
    fill(accent, SOFT_TEAL)
    no_line(accent)
    add_text(slide, 0.68, 0.35, 4.8, 0.35, kicker.upper(), size=12, bold=True, color=TEAL)
    add_text(slide, 11.8, 0.35, 0.9, 0.3, f"{number} / {total}", size=11, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)
    if title:
        add_text(slide, 0.68, 0.88, 11.8, 0.72, title, size=34, bold=True)
    add_text(slide, 0.68, 7.05, 11.8, 0.24, footer, size=10, color=MUTED)
    return slide


def add_table(slide, x, y, w, h, rows, widths=None):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if widths:
        for i, width in enumerate(widths):
            table.columns[i].width = Inches(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(235, 230, 217) if r == 0 else PANEL
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(12 if r == 0 else 13)
                    run.font.bold = r == 0
                    run.font.color.rgb = INK
    return table_shape


slides = [
    ("Developer Webinar", "Developer-Side Memory Optimization for Reducing Token Cost"),
    ("Frame", "This talk is about what our app controls."),
    ("Problem", "Token cost grows when a chatbot keeps sending everything."),
    ("Basics", "Token cost basics"),
    ("Mental Model", "Put memory in the cheapest useful place."),
    ("Technique 1", "Prompt pruning: remove unnecessary prompt text."),
    ("Technique 2", "Conversation summary: compress old chat history."),
    ("Technique 3", "Retrieval/RAG: send relevant document chunks, not the whole library."),
    ("RAG Clarification", "File upload vs custom RAG"),
    ("Technique 4", "Response caching: skip the LLM when an answer already exists."),
    ("Caching Types", "Three practical cache types"),
    ("Technique 5", "Output control: ask for the shortest useful answer."),
    ("Technique 6", "Batch processing: move non-urgent work to the background."),
    ("Technique 6", "Model routing: do not use the strongest model for every task."),
    ("Case Study", "Support chatbot before and after"),
    ("Advanced Names", "Advanced techniques and tools to remember"),
    ("Conclusion", "Before calling the LLM API, ask:"),
]


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    total = len(slides)

    s = base_slide(prs, slides[0][0], "", 1, total, "30-40 minutes | Prepared for August 15, 2026")
    add_text(s, 0.68, 1.45, 11.4, 1.75, slides[0][1], size=42, bold=True)
    add_text(s, 0.72, 3.26, 9.8, 0.5, "For API-based chatbots, websites, support systems, and internal tools.", size=23, color=MUTED)
    x, y = 0.72, 4.15
    for tag in ["Prune", "Summarize", "Retrieve", "Cache answers", "Control output", "Batch", "Route"]:
        used = add_tag(s, x, y, tag)
        x += used + 0.16

    s = base_slide(prs, slides[1][0], slides[1][1], 2, total, "Main correction: we optimize our API system, not the provider's backend.")
    add_card(s, 0.75, 2.0, 5.65, 3.8, "Developer-side decisions", accent=GREEN, bullets=[
        "What history to send", "Which document chunks to send", "When to call the LLM",
        "How long the output should be", "Which model should handle the task"
    ], body_size=17)
    add_card(s, 6.85, 2.0, 5.65, 3.8, "Not the main focus", accent=RED, bullets=[
        "Provider internal caching", "Provider infrastructure tricks", "Exact pricing tables", "Hidden model implementation"
    ], body_size=17)

    s = base_slide(prs, slides[2][0], slides[2][1], 3, total, "Memory is useful, but memory is not free.")
    add_card(s, 0.75, 2.05, 3.72, 2.1, "Input grows", "System prompt, chat history, documents, tools, and retrieved data enter every request.", TEAL, body_size=16)
    add_card(s, 4.82, 2.05, 3.72, 2.1, "Output grows", "Long answers, explanations, and verbose JSON increase generated tokens.", GOLD, body_size=16)
    add_card(s, 8.9, 2.05, 3.72, 2.1, "Calls grow", "Repeated FAQs and offline jobs can trigger unnecessary API calls.", BLUE, body_size=16)
    add_text(s, 1.25, 5.0, 10.9, 0.55, "Cost opportunity = fewer input tokens + fewer output tokens + fewer API calls", size=24, bold=True, align=PP_ALIGN.CENTER)

    s = base_slide(prs, slides[3][0], slides[3][1], 4, total, "Optimization controls what goes in, what comes out, and whether the call is needed.")
    add_card(s, 0.75, 1.9, 5.65, 1.65, "Input tokens", "Instructions, user message, history, documents, tools, examples, and formatting rules.", TEAL, body_size=16)
    add_card(s, 6.85, 1.9, 5.65, 1.65, "Output tokens", "Everything generated by the model. Short useful answers reduce cost and latency.", GOLD, body_size=16)
    add_card(s, 0.75, 4.05, 5.65, 1.65, "Context window", "The maximum information a model can consider in one request.", GREEN, body_size=16)
    add_card(s, 6.85, 4.05, 5.65, 1.65, "API call", "Every model request can cost money. Avoiding a call can be better than shrinking a prompt.", BLUE, body_size=16)

    s = base_slide(prs, slides[4][0], slides[4][1], 5, total, "The model does not know what is important unless our app chooses.")
    for i, (title, body) in enumerate([
        ("Prompt", "Only what the model needs right now."),
        ("Summary", "Compressed old conversation."),
        ("Documents", "Stored outside the prompt and retrieved when needed."),
        ("Answer cache", "Saved answers for repeated questions."),
    ]):
        add_card(s, 0.75 + i * 3.08, 2.1, 2.72, 2.2, title, body, [TEAL, GOLD, GREEN, BLUE][i], body_size=15)
    add_text(s, 1.0, 5.2, 11.3, 0.58, "The app decides what to remove, compress, retrieve, cache, limit, batch, or route.", size=23, bold=True, align=PP_ALIGN.CENTER)

    s = base_slide(prs, slides[5][0], slides[5][1], 6, total, "Names to mention: OpenAI, Claude, DeepSeek, LangChain.")
    add_card(s, 0.85, 2.0, 5.65, 2.1, "Before", "Be helpful. Be friendly. Be concise. Do not be too long. Use simple language. Give bullet points if possible.", RED, body_size=16)
    add_card(s, 6.85, 2.0, 5.65, 2.1, "After", "Answer in simple language using 3 concise bullet points.", GREEN, body_size=18)
    add_table(s, 0.85, 4.75, 11.65, 1.2, [
        ["Manual", "Automatic", "Risk"],
        ["Remove repeated rules, unused examples, vague wording.", "Trim old messages by token limit, last N messages, or relevance.", "Removing important context can hurt quality."]
    ], widths=[3.8, 4.6, 3.25])

    s = base_slide(prs, slides[6][0], slides[6][1], 7, total, "Names to mention: LangChain, OpenAI, Claude, DeepSeek.")
    add_card(s, 0.85, 1.95, 5.65, 1.75, "Before", "Send all 40 previous messages every time the user asks a new question.", RED, body_size=17)
    add_card(s, 6.85, 1.95, 5.65, 1.75, "After", "Send a short summary, important facts, and only the last 4-6 turns.", GREEN, body_size=17)
    add_code(s, 1.45, 4.35, 10.4, 1.45, "Summary:\nUser wants refund for order #123.\nReset email did not arrive.\nUser prefers short step-by-step instructions.", 16)

    s = base_slide(prs, slides[7][0], slides[7][1], 8, total, "Names to mention: Pinecone, Weaviate, Microsoft Azure AI Search, OpenAI vector stores.")
    for i, (title, body) in enumerate([
        ("Store", "PDF, docs, FAQ, policy files."),
        ("Search", "Find chunks matching the question."),
        ("Send", "Only top relevant chunks enter the prompt."),
        ("Answer", "LLM writes using focused context."),
    ]):
        add_card(s, 0.75 + i * 3.08, 2.05, 2.72, 2.15, title, body, [TEAL, GOLD, GREEN, BLUE][i], body_size=15)
    add_text(s, 1.0, 5.1, 11.4, 0.75, "RAG still uses the LLM. It saves input tokens by making the LLM read only the useful part.", size=22, bold=True, align=PP_ALIGN.CENTER)

    s = base_slide(prs, slides[8][0], slides[8][1], 9, total, "RAG saves tokens only if retrieval finds the right chunks.")
    add_card(s, 0.85, 1.95, 5.65, 2.0, "Traditional upload", "A platform may handle document reading internally. Developers usually do not control chunk size, ranking, or exact context.", GOLD, body_size=16)
    add_card(s, 6.85, 1.95, 5.65, 2.0, "Custom RAG chatbot", "Your app stores docs, searches them, selects chunks, and sends only relevant context to the LLM API.", GREEN, body_size=16)
    add_code(s, 0.95, 4.65, 11.4, 1.05, "Without RAG: full_document_tokens + question_tokens + answer_tokens\nWith RAG: search_cost + relevant_chunk_tokens + question_tokens + answer_tokens", 15)

    s = base_slide(prs, slides[9][0], slides[9][1], 10, total, "Names to mention: Redis, Cloudflare AI Gateway, LangChain, vector databases.")
    add_card(s, 0.85, 2.0, 5.65, 2.0, "Without cache", "Same FAQ question -> call LLM -> pay input tokens + output tokens every time.", RED, body_size=16)
    add_card(s, 6.85, 2.0, 5.65, 2.0, "With cache", "Same or similar FAQ question -> return saved answer -> no LLM token cost.", GREEN, body_size=16)
    add_text(s, 1.0, 5.0, 11.3, 0.55, "Best for stable, repeatable questions: refund policy, pricing, password reset, admission rules.", size=22, bold=True, align=PP_ALIGN.CENTER)

    s = base_slide(prs, slides[10][0], slides[10][1], 11, total, "Warning: avoid blind caching for private, sensitive, or fast-changing answers.")
    add_card(s, 0.75, 2.0, 3.72, 1.85, "Exact cache", "Same question text returns the same stored answer.", TEAL, body_size=16)
    add_card(s, 4.82, 2.0, 3.72, 1.85, "FAQ cache", "Approved answers for common support questions.", GREEN, body_size=16)
    add_card(s, 8.9, 2.0, 3.72, 1.85, "Semantic cache", "Similar meaning returns a stored answer if similarity is high enough.", PLUM, body_size=16)
    add_code(s, 2.1, 4.6, 9.0, 1.25, "if similarity(user_question, cached_question) > 0.90:\n  return cached_answer\nelse:\n  call_llm()", 16)

    s = base_slide(prs, slides[11][0], slides[11][1], 12, total, "Names to mention: OpenAI, Claude, DeepSeek, support automation systems.")
    add_card(s, 0.85, 1.95, 5.65, 1.85, "Loose request", "\"Explain everything about password reset.\"", RED, body_size=18)
    add_card(s, 6.85, 1.95, 5.65, 1.85, "Controlled request", "\"Answer in 4 bullets: cause, next step, warning, contact option.\"", GREEN, body_size=17)
    add_table(s, 0.85, 4.45, 11.65, 1.35, [
        ["Developer controls", "Use cases", "Risk"],
        ["max_output_tokens, bullets, JSON, templates, quick/detailed modes.", "Support replies, classifiers, summaries, internal automation.", "Too-low limits can cut off important detail."]
    ], widths=[4.4, 4.1, 3.15])

    s = base_slide(prs, slides[12][0], slides[12][1], 13, total, "Names to mention: OpenAI Batch API, Claude Message Batches, cron jobs, background queues.")
    add_card(s, 0.85, 2.0, 5.65, 2.0, "Bad fit", "User waits while every ticket, summary, or report is processed immediately.", RED, body_size=17)
    add_card(s, 6.85, 2.0, 5.65, 2.0, "Better fit", "Collect, deduplicate, schedule, and process non-urgent LLM jobs later.", GREEN, body_size=17)
    add_bullets(s, 1.1, 4.85, 11.1, 0.9, ["Good for ticket classification, evaluations, summaries, embeddings, reports.", "Not good for live chat answers where the user is waiting."], 17)

    s = base_slide(prs, slides[13][0], slides[13][1], 14, total, "Names to mention: LiteLLM, LangChain, OpenAI, Claude, DeepSeek.")
    for i, (title, body) in enumerate([
        ("No model", "Greeting, cached FAQ, fixed template."),
        ("Cheap model", "Classification, short summary, simple rewrite."),
        ("RAG + model", "Document-specific answers."),
        ("Strong model", "Complex, risky, multi-step reasoning."),
    ]):
        add_card(s, 0.75 + i * 3.08, 2.05, 2.72, 2.15, title, body, [TEAL, GOLD, GREEN, BLUE][i], body_size=15)
    add_text(s, 1.0, 5.15, 11.2, 0.55, "Route by difficulty and risk, not only by cost.", size=24, bold=True, align=PP_ALIGN.CENTER)

    s = base_slide(prs, slides[14][0], slides[14][1], 15, total, "The user still gets support, but the app sends less memory and avoids some calls.")
    add_card(s, 0.85, 1.8, 5.65, 4.1, "Before", accent=RED, bullets=[
        "Strong LLM for every request", "Full policy PDF sent every time", "Full chat history sent every time",
        "Long free-form answers", "Repeated FAQs call the LLM again"
    ], body_size=16)
    add_card(s, 6.85, 1.8, 5.65, 4.1, "After", accent=GREEN, bullets=[
        "Pruned prompt", "Summary plus recent turns", "RAG for policy chunks",
        "Cached FAQ answers", "Short structured replies", "Batch jobs and model routing"
    ], body_size=16)

    s = base_slide(prs, slides[15][0], slides[15][1], 16, total, "Mention these as examples of real developer tools, not mandatory choices.")
    add_table(s, 0.8, 1.75, 11.75, 4.65, [
        ["Technique", "Names to mention"],
        ["RAG with vector databases", "Pinecone, Weaviate, Microsoft Azure AI Search"],
        ["Semantic caching", "Redis, Cloudflare AI Gateway, LangChain"],
        ["LLM routing and fallback", "LiteLLM, LangChain"],
        ["Batch processing", "OpenAI Batch API, Claude Message Batches"],
        ["Support automation patterns", "Zendesk, Intercom"],
    ], widths=[4.45, 7.3])

    s = base_slide(prs, slides[16][0], slides[16][1], 17, total, "End with the framework, then Q&A.")
    x, y = 0.8, 2.15
    for tag in ["Can I remove it?", "Can I summarize it?", "Can I retrieve only part?", "Can I return cached answer?", "Can I limit output?", "Can I batch it?", "Can I route it?"]:
        used = add_tag(s, x, y, tag, w=max(1.65, 0.12 * len(tag) + 0.45))
        x += used + 0.15
        if x > 10.9:
            x = 0.8
            y += 0.58
    add_text(s, 1.0, 5.0, 11.2, 0.95, "Memory optimization is application design: keep the right information in the cheapest useful place.", size=27, bold=True, align=PP_ALIGN.CENTER)

    prs.save(OUT)


if __name__ == "__main__":
    build()
    print(OUT)
